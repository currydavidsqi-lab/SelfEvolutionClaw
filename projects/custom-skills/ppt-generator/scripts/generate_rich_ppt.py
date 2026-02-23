#!/usr/bin/env python3
"""
Enhanced PPT Generator - Rich Content Version
Generate professional PowerPoint with real content
"""

import sys
from pathlib import Path

def create_ai_agent_ppt():
    """Create AI Agent Technology presentation with rich content"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("Error: python-pptx required")
        print("Run: pip install python-pptx")
        sys.exit(1)
    
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Add title manually
    from pptx.util import Pt
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(2)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Agent Technology"
    p.font.size = Pt(54)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(left, Inches(4), width, Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "The Future of Autonomous AI Systems"
    p.font.size = Pt(28)
    p.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(left, Inches(5), width, Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "February 2026"
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: What is an AI Agent?
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "What is an AI Agent?"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    points = [
        "An autonomous system that perceives, reasons, and acts",
        "Can complete complex multi-step tasks without human intervention",
        "Learns from experience and improves over time",
        "Uses tools, memory, and planning to achieve goals",
        "Examples: ChatGPT, Claude, AutoGPT, OpenClaw"
    ]
    
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    # Slide 3: Key Components
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Components of AI Agents"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    components = [
        ("Memory System", "Long-term and short-term context storage"),
        ("Tool Integration", "API calls, file operations, web browsing"),
        ("Planning Engine", "Task decomposition and execution strategy"),
        ("Reasoning Layer", "Chain-of-thought and decision making"),
        ("Learning Module", "Feedback processing and improvement")
    ]
    
    for i, (name, desc) in enumerate(components):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{name}: {desc}"
        p.level = 0
    
    # Slide 4: Popular Frameworks
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Popular AI Agent Frameworks"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    frameworks = [
        "LangChain - 127K stars - Universal LLM framework",
        "AutoGPT - 165K stars - First autonomous agent",
        "MetaGPT - 64K stars - Multi-agent collaboration",
        "DeerFlow - 20K stars - ByteDance SuperAgent",
        "OpenClaw - Personal AI assistant framework"
    ]
    
    for i, fw in enumerate(frameworks):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = fw
        p.level = 0
    
    # Slide 5: Memory Systems
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Memory Systems in AI Agents"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    memory_points = [
        "Short-term: Conversation context (last N messages)",
        "Long-term: Persistent knowledge storage",
        "Episodic: Event-based memories",
        "Semantic: Facts and concepts",
        "Tools: MemOS (5.7K stars), Mem0 (47K stars)"
    ]
    
    for i, point in enumerate(memory_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    # Slide 6: Tool Integration
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Tool Integration: MCP Protocol"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    mcp_points = [
        "Model Context Protocol (MCP) - 81K stars",
        "Standard interface for AI-tool communication",
        "Filesystem, GitHub, Playwright, Database connectors",
        "Supported by Anthropic, OpenAI, Google",
        "Enables 1000+ tool integrations"
    ]
    
    for i, point in enumerate(mcp_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    # Slide 7: Use Cases
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Real-World Applications"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    use_cases = [
        "Software Development: Code generation, testing, debugging",
        "Research: Literature review, data analysis, report writing",
        "Customer Service: 24/7 support, ticket routing, FAQ",
        "Personal Assistant: Scheduling, email, task management",
        "Content Creation: Writing, design, video production"
    ]
    
    for i, case in enumerate(use_cases):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = case
        p.level = 0
    
    # Slide 8: Challenges
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Challenges and Limitations"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    challenges = [
        "Hallucination: Generating false information",
        "Context Length: Limited memory capacity",
        "Cost: High API and compute expenses",
        "Reliability: Inconsistent performance",
        "Safety: Unintended actions and risks"
    ]
    
    for i, challenge in enumerate(challenges):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = challenge
        p.level = 0
    
    # Slide 9: Future Trends
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Future Trends"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    trends = [
        "Multi-Agent Collaboration: Teams of specialized agents",
        "Self-Improvement: Agents that learn and evolve",
        "Embodied AI: Physical robots with AI brains",
        "AGI Progress: Towards general intelligence",
        "Democratization: Everyone can build agents"
    ]
    
    for i, trend in enumerate(trends):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = trend
        p.level = 0
    
    # Slide 10: Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Summary"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    summary = [
        "AI Agents are autonomous systems that can complete complex tasks",
        "Key components: Memory, Tools, Planning, Reasoning, Learning",
        "Major frameworks: LangChain, AutoGPT, MetaGPT, DeerFlow",
        "MCP protocol enables standard tool integration",
        "Challenges remain but progress is rapid",
        "The future: Multi-agent systems and AGI"
    ]
    
    for i, point in enumerate(summary):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    # Save
    output_path = "/tmp/ai-agents-rich.pptx"
    prs.save(output_path)
    print(f"✅ Generated: {output_path}")
    print(f"📊 Slides: 10")
    return output_path

if __name__ == "__main__":
    create_ai_agent_ppt()
