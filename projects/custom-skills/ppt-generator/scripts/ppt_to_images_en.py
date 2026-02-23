#!/usr/bin/env python3
"""
PPT 转图片预览器 - 英文版
将 PPT 每页转换为预览图片（避免中文乱码）
"""

import sys
from pathlib import Path

def ppt_to_images(ppt_path: str, output_dir: str):
    """将 PPT 转换为图片"""
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("错误：需要安装 python-pptx 和 Pillow")
        sys.exit(1)
    
    # 创建输出目录
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    
    # 加载 PPT
    prs = Presentation(ppt_path)
    
    # 图片尺寸
    width, height = 1920, 1080
    bg_color = "#FFFFFF"
    title_color = "#1a1a2e"
    content_color = "#333333"
    accent_color = "#4a90d9"
    
    # 加载字体
    try:
        title_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 64)
        content_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 32)
        small_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
    except:
        title_font = ImageFont.load_default()
        content_font = ImageFont.load_default()
        small_font = ImageFont.load_default()
    
    for i, slide in enumerate(prs.slides, 1):
        # 创建图片
        img = Image.new('RGB', (width, height), bg_color)
        draw = ImageDraw.Draw(img)
        
        # 绘制顶部装饰条
        draw.rectangle([0, 0, width, 100], fill=accent_color)
        
        # 绘制页码
        page_text = f"Slide {i} / {len(prs.slides)}"
        draw.text((width - 180, 35), page_text, fill='white', font=small_font)
        
        # 获取幻灯片内容
        title_text = ""
        content_lines = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if shape == slide.shapes.title:
                    # 翻译常见中文标题
                    translations = {
                        "AI 代理技术": "AI Agent Technology",
                        "目录": "Table of Contents",
                        "总结": "Summary",
                        "第": "Section"
                    }
                    for cn, en in translations.items():
                        text = text.replace(cn, en)
                    title_text = text
                else:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            # 翻译常见中文内容
                            line = para.text.strip()
                            translations = {
                                "副标题": "Subtitle",
                                "作者": "Author",
                                "日期": "Date",
                                "第": "Section",
                                "部分：内容": ": Content",
                                "要点 1：具体内容描述": "Point 1: Detailed description",
                                "要点 2：具体内容描述": "Point 2: Detailed description",
                                "要点 3：具体内容描述": "Point 3: Detailed description",
                                "核心要点回顾": "Key points review",
                                "关键结论": "Key conclusions",
                                "后续行动": "Next steps"
                            }
                            for cn, en in translations.items():
                                line = line.replace(cn, en)
                            content_lines.append(line)
        
        # 绘制标题
        if title_text:
            y_pos = 180
            draw.text((80, y_pos), title_text, fill=title_color, font=title_font)
        
        # 绘制内容
        y_pos = 320
        for line in content_lines[:6]:  # 最多显示6行
            # 绘制项目符号
            draw.ellipse([90, y_pos + 12, 110, y_pos + 32], fill=accent_color)
            draw.text((130, y_pos), line, fill=content_color, font=content_font)
            y_pos += 70
        
        # 保存图片
        img_path = output_path / f"slide_{i:02d}.png"
        img.save(img_path, 'PNG')
        print(f"✅ Generated: {img_path}")
    
    print(f"\n📊 Total: {len(prs.slides)} slides")
    print(f"📁 Location: {output_dir}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python ppt_to_images.py <ppt_file> <output_dir>")
        print("Example: python ppt_to_images.py /tmp/demo.pptx /tmp/slides")
        sys.exit(1)
    
    ppt_to_images(sys.argv[1], sys.argv[2])
