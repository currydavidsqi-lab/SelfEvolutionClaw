#!/usr/bin/env python3
"""
PPT 转图片预览器
将 PPT 每页转换为预览图片
"""

import sys
from pathlib import Path

def ppt_to_images(ppt_path: str, output_dir: str):
    """将 PPT 转换为图片"""
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("错误：需要安装 python-pptx 和 Pillow")
        sys.exit(1)
    
    # 创建输出目录
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    
    # 加载 PPT
    prs = Presentation(ppt_path)
    
    # 图片尺寸
    width, height = 1920, 1080
    bg_color = "#FFFFFF"
    title_color = "#1a1a2e"
    content_color = "#333333"
    accent_color = "#4a90d9"
    
    # 尝试加载字体
    try:
        title_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 72)
        content_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 36)
    except:
        title_font = ImageFont.load_default()
        content_font = ImageFont.load_default()
    
    for i, slide in enumerate(prs.slides, 1):
        # 创建图片
        img = Image.new('RGB', (width, height), bg_color)
        draw = ImageDraw.Draw(img)
        
        # 绘制顶部装饰条
        draw.rectangle([0, 0, width, 120], fill=accent_color)
        
        # 绘制页码
        page_text = f"Slide {i}/{len(prs.slides)}"
        draw.text((width - 200, 40), page_text, fill='white', font=content_font)
        
        # 获取幻灯片内容
        title_text = ""
        content_lines = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if shape == slide.shapes.title:
                    title_text = text
                else:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            content_lines.append(para.text.strip())
        
        # 绘制标题
        if title_text:
            y_pos = 200
            draw.text((100, y_pos), title_text, fill=title_color, font=title_font)
        
        # 绘制内容
        y_pos = 350
        for line in content_lines[:6]:  # 最多显示6行
            # 绘制项目符号
            draw.ellipse([110, y_pos + 15, 130, y_pos + 35], fill=accent_color)
            draw.text((150, y_pos), line, fill=content_color, font=content_font)
            y_pos += 80
        
        # 保存图片
        img_path = output_path / f"slide_{i:02d}.png"
        img.save(img_path, 'PNG')
        print(f"✅ 生成：{img_path}")
    
    print(f"\n📊 共生成 {len(prs.slides)} 张图片")
    print(f"📁 保存位置：{output_dir}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("用法：python ppt_to_images.py <ppt文件> <输出目录>")
        print("示例：python ppt_to_images.py /tmp/demo.pptx /tmp/slides")
        sys.exit(1)
    
    ppt_to_images(sys.argv[1], sys.argv[2])
