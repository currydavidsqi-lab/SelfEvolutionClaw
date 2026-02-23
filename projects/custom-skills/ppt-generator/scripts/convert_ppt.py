#!/usr/bin/env python3
"""
PPT to PNG Converter - Enhanced Version
Convert PowerPoint slides to high-quality PNG images
"""

import sys
from pathlib import Path

def convert_ppt_to_png(ppt_path: str, output_dir: str):
    """Convert PPT slides to PNG images"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("Error: python-pptx and Pillow required")
        sys.exit(1)
    
    # Create output directory
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    
    # Load PPT
    prs = Presentation(ppt_path)
    
    # Image dimensions (1080p)
    width, height = 1920, 1080
    bg_color = "#FFFFFF"
    title_bg = "#2563EB"
    title_color = "#1E3A8A"
    content_color = "#374151"
    
    # Load fonts
    try:
        title_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 52)
        content_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 28)
        small_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 22)
        large_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 72)
    except:
        title_font = ImageFont.load_default()
        content_font = ImageFont.load_default()
        small_font = ImageFont.load_default()
        large_font = ImageFont.load_default()
    
    for i, slide in enumerate(prs.slides, 1):
        # Create image
        img = Image.new('RGB', (width, height), bg_color)
        draw = ImageDraw.Draw(img)
        
        # Extract text content
        title_text = ""
        content_lines = []
        is_title_slide = (i == 1)  # First slide is title
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    
                    # Check if it's a title
                    if shape == slide.shapes.title or (is_title_slide and len(text) < 50):
                        if not title_text:
                            title_text = text
                            continue
                    
                    # It's content
                    content_lines.append(text)
        
        # Draw title slide differently
        if is_title_slide:
            # Title background
            draw.rectangle([0, 0, width, height], fill=title_bg)
            
            # Main title
            if title_text:
                bbox = draw.textbbox((0, 0), title_text, font=large_font)
                text_width = bbox[2] - bbox[0]
                x = (width - text_width) // 2
                draw.text((x, 350), title_text, fill='white', font=large_font)
            
            # Subtitle/content
            y_pos = 500
            for line in content_lines[:2]:  # First 2 lines only
                bbox = draw.textbbox((0, 0), line, font=content_font)
                text_width = bbox[2] - bbox[0]
                x = (width - text_width) // 2
                draw.text((x, y_pos), line, fill='#E0E7FF', font=content_font)
                y_pos += 60
        
        else:
            # Regular slide
            # Header bar
            draw.rectangle([0, 0, width, 120], fill=title_bg)
            
            # Slide number
            page_text = f"{i} / {len(prs.slides)}"
            bbox = draw.textbbox((0, 0), page_text, font=small_font)
            text_width = bbox[2] - bbox[0]
            draw.text((width - text_width - 40, 45), page_text, fill='white', font=small_font)
            
            # Title
            if title_text:
                draw.text((60, 25), title_text, fill='white', font=title_font)
            
            # Content
            y_pos = 180
            for line in content_lines[:8]:  # Max 8 lines
                # Bullet point
                draw.ellipse([70, y_pos + 10, 88, y_pos + 28], fill='#3B82F6')
                
                # Text (wrap if too long)
                if len(line) > 90:
                    # Split into two lines
                    words = line.split()
                    mid = len(words) // 2
                    line1 = ' '.join(words[:mid])
                    line2 = ' '.join(words[mid:])
                    draw.text((100, y_pos), line1, fill=content_color, font=content_font)
                    y_pos += 50
                    draw.text((100, y_pos), line2, fill=content_color, font=content_font)
                else:
                    draw.text((100, y_pos), line, fill=content_color, font=content_font)
                
                y_pos += 70
        
        # Save image
        img_path = output_path / f"slide_{i:02d}.png"
        img.save(img_path, 'PNG', quality=95)
        print(f"✅ Slide {i}: {img_path}")
    
    print(f"\n✨ Total: {len(prs.slides)} slides converted")
    print(f"📁 Location: {output_dir}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python convert_ppt.py <ppt_file> <output_dir>")
        sys.exit(1)
    
    convert_ppt_to_png(sys.argv[1], sys.argv[2])
