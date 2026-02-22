#!/usr/bin/env python3
"""
PPT 生成器 v0.1
根据主题自动生成 PowerPoint 演示文稿
"""

import argparse
import sys
from pathlib import Path

def generate_outline(topic: str, slides: int) -> list:
    """生成 PPT 大纲"""
    outline = []
    
    # 第一张：标题页
    outline.append({
        "title": topic,
        "content": ["副标题", "作者", "日期"]
    })
    
    # 第二张：目录
    outline.append({
        "title": "目录",
        "content": [f"第 {i+1} 部分：内容 {i+1}" for i in range(min(5, slides-2))]
    })
    
    # 中间页：内容
    for i in range(2, slides-1):
        outline.append({
            "title": f"第 {i-1} 部分：内容",
            "content": [
                "要点 1：具体内容描述",
                "要点 2：具体内容描述",
                "要点 3：具体内容描述"
            ]
        })
    
    # 最后一张：总结
    if slides > 2:
        outline.append({
            "title": "总结",
            "content": [
                "核心要点回顾",
                "关键结论",
                "后续行动"
            ]
        })
    
    return outline

def create_pptx(outline: list, output_path: str):
    """创建 PPTX 文件"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("错误：需要安装 python-pptx")
        print("运行：pip install python-pptx")
        sys.exit(1)
    
    prs = Presentation()
    
    for slide_data in outline:
        # 使用标题+内容布局
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = slide_data["title"]
        
        # 设置内容
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        for i, point in enumerate(slide_data["content"]):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = point
            p.level = 0
    
    # 保存文件
    prs.save(output_path)
    print(f"✅ PPT 已生成：{output_path}")
    print(f"📊 幻灯片数量：{len(outline)}")

def main():
    parser = argparse.ArgumentParser(
        description="PPT 生成器 - 根据主题自动生成演示文稿",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例：
  %(prog)s --topic "AI 代理技术" --slides 10
  %(prog)s --topic "机器学习基础" --slides 8 --output presentation.pptx
        """
    )
    
    parser.add_argument(
        "--topic",
        required=True,
        help="PPT 主题"
    )
    parser.add_argument(
        "--slides",
        type=int,
        default=10,
        help="幻灯片数量（默认：10）"
    )
    parser.add_argument(
        "--output",
        default="output.pptx",
        help="输出文件路径（默认：output.pptx）"
    )
    
    args = parser.parse_args()
    
    # 生成大纲
    print(f"📝 生成主题：{args.topic}")
    print(f"📊 幻灯片数：{args.slides}")
    outline = generate_outline(args.topic, args.slides)
    
    # 创建 PPT
    create_pptx(outline, args.output)

if __name__ == "__main__":
    main()
