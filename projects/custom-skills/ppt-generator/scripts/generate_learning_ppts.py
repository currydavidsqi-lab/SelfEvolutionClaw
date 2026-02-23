#!/usr/bin/env python3
"""
Generate Learning Summary PPTs
分类整理学习内容成多个 PPT
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN

def create_ppt(title: str, slides_data: list, output_path: str):
    """创建 PPT"""
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title background
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    
    # Add title text box
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(2)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(left, Inches(4.5), width, Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CurryManager Learning Summary"
    p.font.size = Pt(28)
    p.alignment = PP_ALIGN.CENTER
    
    # Content slides
    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = slide_data['title']
        
        # Content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        for i, point in enumerate(slide_data['content']):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.level = 0
    
    prs.save(output_path)
    print(f"✅ Generated: {output_path}")
    return output_path


def ppt1_ai_frameworks():
    """PPT 1: AI Agent Frameworks"""
    slides = [
        {
            "title": "Top AI Agent Frameworks (2026)",
            "content": [
                "LangChain - 127K stars - Universal LLM framework",
                "AutoGPT - 165K stars - First autonomous agent",
                "MetaGPT - 64K stars - Multi-agent collaboration",
                "DeerFlow - 20K stars - ByteDance SuperAgent",
                "Letta - 21K stars - Stateful agent platform"
            ]
        },
        {
            "title": "LangChain - Most Popular",
            "content": [
                "Universal framework for LLM applications",
                "Components: Chains, Agents, Memory, Tools",
                "Supports 50+ LLM providers",
                "Production-ready with monitoring",
                "Active community: 2000+ contributors"
            ]
        },
        {
            "title": "AutoGPT - Autonomous Pioneer",
            "content": [
                "First fully autonomous AI agent",
                "Self-prompting and task planning",
                "Internet access and file operations",
                "Inspired thousands of clones",
                "Still evolving: GPT-Engineer, AgentGPT"
            ]
        },
        {
            "title": "MetaGPT - Multi-Agent Teams",
            "content": [
                "Software company simulation",
                "Roles: Product Manager, Architect, Engineer",
                "Complete software from one prompt",
                "Collaborative problem solving",
                "Used by major tech companies"
            ]
        },
        {
            "title": "DeerFlow - ByteDance SuperAgent",
            "content": [
                "Handle tasks from minutes to hours",
                "Skills, Tools, Sub-agents support",
                "Sandbox and file system",
                "Context engineering built-in",
                "Long-term memory integration"
            ]
        }
    ]
    
    create_ppt(
        "AI Agent Frameworks",
        slides,
        "/tmp/ppt1_frameworks.pptx"
    )


def ppt2_automation_platforms():
    """PPT 2: AI Automation Platforms"""
    slides = [
        {
            "title": "Top AI Automation Platforms (2026)",
            "content": [
                "n8n - 176K stars - Workflow automation",
                "Dify - 130K stars - LLM app builder",
                "Apache Airflow - 38K stars - Data pipeline",
                "Activepieces - 15K stars - Open-source Zapier"
            ]
        },
        {
            "title": "n8n - Workflow Automation Leader",
            "content": [
                "400+ integrations available",
                "Visual drag-and-drop builder",
                "Self-hosted or cloud option",
                "JavaScript customization",
                "Active community: 500+ templates"
            ]
        },
        {
            "title": "Dify - LLM App Builder",
            "content": [
                "ByteDance open-source project",
                "No-code LLM application builder",
                "RAG, Agent, Workflow support",
                "Multi-model support",
                "Enterprise-ready deployment"
            ]
        },
        {
            "title": "Key Features Comparison",
            "content": [
                "n8n: Best for integrations, visual workflow",
                "Dify: Best for LLM apps, no-code",
                "Airflow: Best for data pipelines, scheduling",
                "Activepieces: Best for simple automations",
                "All support self-hosting"
            ]
        },
        {
            "title": "Integration with OpenClaw",
            "content": [
                "n8n: Webhook triggers, API calls",
                "Dify: LLM integration patterns",
                "Airflow: Scheduled tasks",
                "Activepieces: Quick automations",
                "Choose based on use case"
            ]
        }
    ]
    
    create_ppt(
        "AI Automation Platforms",
        slides,
        "/tmp/ppt2_automation.pptx"
    )


def ppt3_llm_training():
    """PPT 3: LLM Training & Fine-tuning"""
    slides = [
        {
            "title": "LLM Training & Fine-tuning Ecosystem",
            "content": [
                "LlamaFactory - 67K stars - Unified fine-tuning",
                "Llama Cookbook - 18K stars - Meta official guide",
                "Axolotl - 9K stars - Fine-tuning toolkit",
                "LLaMA-Factory - Production ready",
                "Unsloth - Fast fine-tuning"
            ]
        },
        {
            "title": "LlamaFactory - Universal Fine-tuner",
            "content": [
                "Fine-tune 100+ LLMs",
                "LoRA, QLoRA, Full fine-tuning",
                "Web UI + CLI interface",
                "Supports all major models",
                "Production deployment ready"
            ]
        },
        {
            "title": "Popular Fine-tuning Methods",
            "content": [
                "Full Fine-tuning: All parameters updated",
                "LoRA: Low-rank adaptation, efficient",
                "QLoRA: Quantized LoRA, memory-saving",
                "Prefix Tuning: Task-specific prefixes",
                "Instruction Tuning: Follow instructions"
            ]
        },
        {
            "title": "Evaluation & Prompt Engineering",
            "content": [
                "Prompt Engineering Guide - 71K stars",
                "OpenAI Evals - Official evaluation framework",
                "OpenCompass - Shanghai AI Lab",
                "Langfuse - LLM observability",
                "Prompt metrics: Quality, Safety, Cost"
            ]
        },
        {
            "title": "LLM Inference & Serving",
            "content": [
                "vLLM - 71K stars - Fastest inference",
                "SGLang - Structured generation",
                "LightLLM - Lightweight serving",
                "Mooncake - KVCache optimization",
                "All support production deployment"
            ]
        }
    ]
    
    create_ppt(
        "LLM Training & Serving",
        slides,
        "/tmp/ppt3_training.pptx"
    )


def ppt4_rag_memory():
    """PPT 4: RAG & Memory Systems"""
    slides = [
        {
            "title": "RAG Technology Ecosystem",
            "content": [
                "RAGFlow - 73K stars - Popular RAG engine",
                "GraphRAG - 31K stars - Microsoft knowledge graph",
                "LightRAG - 29K stars - Simple and fast",
                "RAG Techniques - Comprehensive guide",
                "RAG + Agent = Standard architecture"
            ]
        },
        {
            "title": "RAGFlow - Leading RAG Engine",
            "content": [
                "Document parsing and indexing",
                "Multi-modal support",
                "Graph-based retrieval",
                "Enterprise deployment",
                "Active development and support"
            ]
        },
        {
            "title": "GraphRAG - Microsoft Innovation",
            "content": [
                "Knowledge graph integration",
                "Better reasoning capabilities",
                "Entity relationship extraction",
                "Hierarchical summarization",
                "Research-grade implementation"
            ]
        },
        {
            "title": "AI Memory Systems",
            "content": [
                "MemOS - 5.7K stars - OpenClaw support",
                "Mem0 - 47K stars - Universal memory layer",
                "Letta - 21K stars - Stateful agents",
                "Memory types: Short-term, Long-term, Episodic",
                "Essential for continuous AI agents"
            ]
        },
        {
            "title": "MCP Protocol - Tool Integration",
            "content": [
                "81K stars - Model Context Protocol",
                "Standard AI-tool communication",
                "Filesystem, GitHub, Playwright connectors",
                "Supported by Anthropic, OpenAI, Google",
                "Enables 1000+ tool integrations"
            ]
        }
    ]
    
    create_ppt(
        "RAG & Memory Systems",
        slides,
        "/tmp/ppt4_rag_memory.pptx"
    )


if __name__ == "__main__":
    print("Generating Learning Summary PPTs...")
    print("=" * 50)
    
    ppt1_ai_frameworks()
    ppt2_automation_platforms()
    ppt3_llm_training()
    ppt4_rag_memory()
    
    print("=" * 50)
    print("✅ All PPTs generated!")
    print("\nFiles:")
    print("  1. /tmp/ppt1_frameworks.pptx")
    print("  2. /tmp/ppt2_automation.pptx")
    print("  3. /tmp/ppt3_training.pptx")
    print("  4. /tmp/ppt4_rag_memory.pptx")
