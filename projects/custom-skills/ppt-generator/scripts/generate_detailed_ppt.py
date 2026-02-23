#!/usr/bin/env python3
"""
Enhanced PPT Generator with Charts and Rich Content
包含图表、详细内容的专业 PPT
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title: str, subtitle: str):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景色
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(37, 99, 235)
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(219, 234, 254)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title: str, content: list):
    """添加内容页（详细列表）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(37, 99, 235)
    header.line.fill.background()
    
    # 标题文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 内容区域
    y_pos = 1.6
    for item in content:
        # 项目符号（圆点）
        bullet = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.6), Inches(y_pos + 0.1), Inches(0.15), Inches(0.15)
        )
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = RGBColor(59, 130, 246)
        bullet.line.fill.background()
        
        # 文字
        text_box = slide.shapes.add_textbox(Inches(0.9), Inches(y_pos), Inches(8.5), Inches(0.6))
        tf = text_box.text_frame
        tf.word_wrap = True
        
        if isinstance(item, dict):
            # 标题 + 描述格式
            p = tf.paragraphs[0]
            p.text = item['title']
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(30, 58, 138)
            
            if 'desc' in item:
                p = tf.add_paragraph()
                p.text = item['desc']
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(107, 114, 128)
        else:
            # 纯文字
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(55, 65, 81)
        
        y_pos += 0.8
    
    return slide

def add_chart_slide(prs, title: str, data: list):
    """添加柱状图页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(37, 99, 235)
    header.line.fill.background()
    
    # 标题文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 绘制柱状图
    max_val = max(item['value'] for item in data)
    chart_left = 1.5
    chart_bottom = 6.5
    chart_width = 7
    chart_height = 4.5
    bar_width = chart_width / len(data) * 0.7
    bar_gap = chart_width / len(data) * 0.3
    
    for i, item in enumerate(data):
        # 计算柱子高度
        bar_height = (item['value'] / max_val) * chart_height
        
        # 柱子位置
        left = chart_left + i * (bar_width + bar_gap)
        top = chart_bottom - bar_height
        
        # 绘制柱子
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top),
            Inches(bar_width), Inches(bar_height)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(59, 130, 246)
        bar.line.color.rgb = RGBColor(37, 99, 235)
        
        # 数值标签
        val_box = slide.shapes.add_textbox(
            Inches(left), Inches(top - 0.4),
            Inches(bar_width), Inches(0.4)
        )
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{item['value']:.0f}K"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(37, 99, 235)
        p.alignment = PP_ALIGN.CENTER
        
        # 名称标签
        name_box = slide.shapes.add_textbox(
            Inches(left - 0.2), Inches(chart_bottom + 0.1),
            Inches(bar_width + 0.4), Inches(0.5)
        )
        tf = name_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item['name']
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(55, 65, 81)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_comparison_slide(prs, title: str, items: list):
    """添加对比表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(37, 99, 235)
    header.line.fill.background()
    
    # 标题文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 表格
    y_pos = 1.8
    for i, item in enumerate(items):
        # 行背景（交替颜色）
        row_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(y_pos),
            Inches(9), Inches(0.9)
        )
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = RGBColor(249, 250, 251) if i % 2 == 0 else RGBColor(255, 255, 255)
        row_bg.line.fill.background()
        
        # 名称
        name_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(2.5), Inches(0.7))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = item['name']
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(30, 58, 138)
        
        # Stars
        stars_box = slide.shapes.add_textbox(Inches(3.3), Inches(y_pos + 0.2), Inches(1.2), Inches(0.5))
        tf = stars_box.text_frame
        p = tf.paragraphs[0]
        p.text = item['stars']
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(234, 179, 8)
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(4.6), Inches(y_pos + 0.1), Inches(4.7), Inches(0.7))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item['desc']
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(107, 114, 128)
        
        y_pos += 0.9
    
    return slide

def create_frameworks_ppt():
    """创建 AI Frameworks 详细版 PPT"""
    prs = Presentation()
    
    # Slide 1: 标题
    add_title_slide(prs, 
        "AI Agent Frameworks 2026",
        "Comprehensive Guide to Top Frameworks • CurryManager Learning Summary"
    )
    
    # Slide 2: 概览图表
    add_chart_slide(prs, "GitHub Stars Comparison (K)", [
        {"name": "AutoGPT", "value": 165},
        {"name": "LangChain", "value": 127},
        {"name": "MetaGPT", "value": 64},
        {"name": "Letta", "value": 21},
        {"name": "DeerFlow", "value": 20},
    ])
    
    # Slide 3: 对比表格
    add_comparison_slide(prs, "Framework Comparison", [
        {"name": "LangChain", "stars": "⭐ 127K", "desc": "Universal framework, 50+ LLM providers, production-ready"},
        {"name": "AutoGPT", "stars": "⭐ 165K", "desc": "First autonomous agent, self-prompting, internet access"},
        {"name": "MetaGPT", "stars": "⭐ 64K", "desc": "Multi-agent teams, software company simulation"},
        {"name": "DeerFlow", "stars": "⭐ 20K", "desc": "ByteDance SuperAgent, long tasks, sandbox support"},
        {"name": "Letta", "stars": "⭐ 21K", "desc": "Stateful agents, advanced memory, self-improvement"},
    ])
    
    # Slide 4: LangChain 详细
    add_content_slide(prs, "LangChain - The Universal Framework", [
        {"title": "📦 Core Components", "desc": "Chains, Agents, Memory, Tools, Retrievers - modular building blocks"},
        {"title": "🔌 LLM Support", "desc": "50+ providers: OpenAI, Anthropic, Google, HuggingFace, local models"},
        {"title": "🏭 Production Ready", "desc": "LangSmith for monitoring, LangServe for deployment, evaluation tools"},
        {"title": "👥 Community", "desc": "2000+ contributors, 100+ templates, extensive documentation"},
        {"title": "💡 Best For", "desc": "Building custom LLM applications, RAG systems, chatbots, agents"},
        {"title": "⚠️ Limitations", "desc": "Learning curve, complexity for simple tasks, token costs"},
    ])
    
    # Slide 5: AutoGPT 详细
    add_content_slide(prs, "AutoGPT - Autonomous Pioneer", [
        {"title": "🤖 Core Concept", "desc": "Fully autonomous agent that self-prompts and plans tasks"},
        {"title": "🌐 Capabilities", "desc": "Internet access, file operations, code execution, API calls"},
        {"title": "🎯 How It Works", "desc": "Goal → Sub-tasks → Execute → Evaluate → Iterate"},
        {"title": "🌟 Impact", "desc": "Inspired 1000+ clones: GPT-Engineer, AgentGPT, BabyAGI"},
        {"title": "💡 Best For", "desc": "Research, experimentation, autonomous task completion"},
        {"title": "⚠️ Challenges", "desc": "Can get stuck in loops, expensive API calls, reliability issues"},
    ])
    
    # Slide 6: MetaGPT 详细
    add_content_slide(prs, "MetaGPT - Multi-Agent Collaboration", [
        {"title": "👥 Team Simulation", "desc": "Simulates a software company with specialized agents"},
        {"title": "🎭 Roles", "desc": "Product Manager, Architect, Engineer, QA, Project Manager"},
        {"title": "🔄 Workflow", "desc": "Requirements → Design → Code → Test → Deploy"},
        {"title": "✨ Unique Feature", "desc": "Complete software project from one prompt"},
        {"title": "💼 Use Cases", "desc": "Rapid prototyping, MVP development, education"},
        {"title": "📊 Results", "desc": "Generated games, apps, tools in minutes"},
    ])
    
    # Slide 7: DeerFlow 详细
    add_content_slide(prs, "DeerFlow - ByteDance SuperAgent", [
        {"title": "⏱️ Long-Running Tasks", "desc": "Handle tasks from minutes to hours, not just quick Q&A"},
        {"title": "🛠️ Components", "desc": "Skills, Tools, Sub-agents, Sandbox, File System"},
        {"title": "🧠 Context Engineering", "desc": "Advanced context management for complex tasks"},
        {"title": "💾 Memory", "desc": "Long-term memory integration, cross-session persistence"},
        {"title": "🔒 Safety", "desc": "Sandbox execution, permission control, audit logs"},
        {"title": "🏢 Backed By", "desc": "ByteDance (TikTok parent), enterprise-grade"},
    ])
    
    # Slide 8: 选择指南
    add_content_slide(prs, "Framework Selection Guide", [
        {"title": "🎓 Learning & Prototyping", "desc": "LangChain - Best documentation, large community"},
        {"title": "🚀 Quick Autonomous Tasks", "desc": "AutoGPT - Simple setup, self-driven"},
        {"title": "🏢 Team Projects", "desc": "MetaGPT - Role-based collaboration"},
        {"title": "⏰ Long-Running Tasks", "desc": "DeerFlow - Hours-long tasks, enterprise features"},
        {"title": "🧠 Memory-Intensive Apps", "desc": "Letta - Advanced memory, self-improvement"},
        {"title": "💡 Recommendation", "desc": "Start with LangChain, explore others based on needs"},
    ])
    
    # 保存
    output_path = "/tmp/frameworks_detailed.pptx"
    prs.save(output_path)
    print(f"✅ Generated: {output_path}")
    print(f"📊 Slides: 8")
    return output_path

if __name__ == "__main__":
    print("Creating Enhanced PPT with Charts...")
    print("=" * 50)
    create_frameworks_ppt()
    print("=" * 50)
    print("✨ Done!")
